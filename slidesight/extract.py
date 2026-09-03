"""Pull images and their slide context out of a .pptx.

Two facts drive this module, both verified against real decks:

1. Images hide inside groups. Iterating ``slide.shapes`` misses them; you have
   to recurse into ``MSO_SHAPE_TYPE.GROUP``.
2. Alt text lives in the OOXML ``descr`` attribute on ``cNvPr``, and pictures
   also carry a ``title`` attribute holding the original filename.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class SlideContext:
    """The text surrounding an image, from its own slide only.

    Deliberately per-slide. Sending the whole deck's text with every image
    dilutes the signal -- when describing slide 7's chart, slide 32's bullets
    are noise.
    """

    number: int
    title: str = ""
    body: str = ""
    notes: str = ""


@dataclass
class ImageRef:
    """One picture, everything needed to describe it and write back to it."""

    slide: int
    image_id: str
    blob: bytes
    content_type: str
    context: SlideContext
    existing_alt: str = ""
    existing_title: str = ""
    shape: object = field(default=None, repr=False)


def walk(shapes) -> Iterator:
    """Yield every shape, recursing into groups."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(shape.shapes)


def _cnvpr(shape):
    """The OOXML non-visual properties element that carries alt text."""
    return shape._element._nvXxPr.cNvPr


def get_alt_text(shape) -> str:
    return _cnvpr(shape).attrib.get("descr", "")


def set_alt_text(shape, text: str) -> None:
    """Write alt text and strip the filename ``title``.

    lxml handles escaping, so ampersands and non-ASCII round-trip correctly.
    Never hand-escape here.
    """
    attrib = _cnvpr(shape).attrib
    attrib["descr"] = text
    attrib.pop("title", None)


def slide_context(slide, number: int) -> SlideContext:
    """Title, body text, and speaker notes for one slide.

    No text cleaning. The raw fragmented python-pptx output goes to the model
    as-is; regex stitching is a time sink that buys nothing.
    """
    title = ""
    if slide.shapes.title is not None:
        title = (slide.shapes.title.text or "").strip()

    body_parts: list[str] = []
    for shape in walk(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if text and text != title:
            body_parts.append(text)

    notes = ""
    if slide.has_notes_slide:
        # A notes slide can exist with no notes placeholder on it, in which
        # case notes_text_frame is None rather than an empty frame.
        frame = slide.notes_slide.notes_text_frame
        if frame is not None:
            notes = (frame.text or "").strip()

    return SlideContext(
        number=number,
        title=title,
        body="\n".join(body_parts),
        notes=notes,
    )


def extract_images(prs: Presentation) -> list[ImageRef]:
    """Every picture in the deck, in slide order, with its slide's context."""
    images: list[ImageRef] = []
    for index, slide in enumerate(prs.slides, start=1):
        context = slide_context(slide, index)
        for shape in walk(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            attrib = _cnvpr(shape).attrib
            try:
                blob = shape.image.blob
                content_type = shape.image.content_type
            except Exception:
                # Linked-but-not-embedded pictures have no blob to read.
                continue
            images.append(
                ImageRef(
                    slide=index,
                    image_id=f"shape_{attrib.get('id', len(images))}",
                    blob=blob,
                    content_type=content_type,
                    context=context,
                    existing_alt=attrib.get("descr", ""),
                    existing_title=attrib.get("title", ""),
                    shape=shape,
                )
            )
    return images


def deck_text(prs: Presentation, limit: int = 12000) -> str:
    """All slide text, for the one-per-deck summary call."""
    parts: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        context = slide_context(slide, index)
        chunk = " ".join(x for x in (context.title, context.body) if x)
        if chunk:
            parts.append(f"[Slide {index}] {chunk}")
    return "\n".join(parts)[:limit]


def count_tables(prs: Presentation) -> int:
    """Tables are detected and reported, never remediated."""
    return sum(
        1
        for slide in prs.slides
        for shape in walk(slide.shapes)
        if getattr(shape, "has_table", False)
    )
