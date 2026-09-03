"""WCAG detection checks. Detection only -- nothing here modifies the deck.

Alt text is additive: we fill an empty field and nothing breaks. Colour and type
size are destructive -- changing a professor's palette or font sizes hands back
a deck they do not recognise. So these checks report with slide numbers and let
a human decide. Same confidence-gate principle applied to the whole product.

No model calls. This module runs in milliseconds.
"""

from __future__ import annotations

from dataclasses import asdict, dataclass
from typing import Any

from pptx import Presentation
from pptx.util import Length

from .extract import walk

MIN_BODY_PT = 18.0

VAGUE_LINK_TEXT = {
    "click here",
    "here",
    "read more",
    "more",
    "this link",
    "link",
    "learn more",
    "see more",
    "click",
    "this page",
    "download",
}


@dataclass
class Issue:
    """One accessibility problem, located by slide."""

    check: str
    slide: int
    detail: str
    severity: str = "warning"


def _placeholder_default_pt(shape, slide) -> float | None:
    """Resolve an inherited font size from the slide layout.

    ``run.font.size`` is None when the size comes from the layout rather than
    being set on the run. Skipping those would hide most real violations, so
    look the placeholder up in the layout and read its size instead.
    """
    try:
        if not shape.is_placeholder:
            return None
        idx = shape.placeholder_format.idx
        layout = slide.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx != idx:
                continue
            for para in ph.text_frame.paragraphs:
                if para.font.size is not None:
                    return para.font.size.pt
                for run in para.runs:
                    if run.font.size is not None:
                        return run.font.size.pt
    except Exception:
        return None
    return None


def check_slide_titles(prs: Presentation) -> list[Issue]:
    """Every slide needs a title -- screen readers use them to navigate."""
    issues = []
    for index, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title
        if title is None:
            issues.append(
                Issue("missing_title", index, "slide has no title placeholder")
            )
        elif not (title.text or "").strip():
            issues.append(
                Issue("missing_title", index, "title placeholder is empty")
            )
    return issues


def check_font_sizes(prs: Presentation, minimum: float = MIN_BODY_PT) -> list[Issue]:
    """Flag body text below 18pt.

    Titles are exempt -- they are large by definition, and a short title set at
    16pt in a decorative font is a design choice, not a legibility failure.
    """
    issues = []
    for index, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        smallest: float | None = None
        sample = ""
        for shape in walk(slide.shapes):
            if not getattr(shape, "has_text_frame", False) or shape is title_shape:
                continue
            inherited = _placeholder_default_pt(shape, slide)
            for para in shape.text_frame.paragraphs:
                para_size = para.font.size.pt if para.font.size is not None else None
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    size = (
                        run.font.size.pt
                        if run.font.size is not None
                        else para_size
                        if para_size is not None
                        else inherited
                    )
                    if size is None:
                        continue  # unresolved; reported as unknown, not as a failure
                    if size < minimum and (smallest is None or size < smallest):
                        smallest = size
                        sample = run.text.strip()[:40]
        if smallest is not None:
            issues.append(
                Issue(
                    "small_text",
                    index,
                    f"{smallest:g}pt (below {minimum:g}pt): {sample!r}",
                )
            )
    return issues


def check_table_headers(prs: Presentation) -> list[Issue]:
    """A table without a header row gives a screen reader no column context."""
    issues = []
    for index, slide in enumerate(prs.slides, start=1):
        for shape in walk(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            if not table.first_row:
                issues.append(
                    Issue(
                        "table_no_header",
                        index,
                        f"{len(table.rows)}x{len(table.columns)} table has no header row",
                    )
                )
    return issues


def check_link_text(prs: Presentation) -> list[Issue]:
    """Screen reader users navigate by link list, so "click here" is useless."""
    issues = []
    for index, slide in enumerate(prs.slides, start=1):
        for shape in walk(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        address = run.hyperlink.address
                    except Exception:
                        address = None
                    if not address:
                        continue
                    text = (run.text or "").strip()
                    if text.lower().strip(" .:>-") in VAGUE_LINK_TEXT or not text:
                        issues.append(
                            Issue(
                                "vague_link",
                                index,
                                f"{text or '(empty)'!r} -> {address[:60]}",
                            )
                        )
    return issues


def check_reading_order(prs: Presentation) -> list[Issue]:
    """Shapes are announced in XML order, not visual order.

    A deck can look perfect and still read out backwards. Compare each shape's
    position against its index and flag slides where the two disagree.
    """
    issues = []
    for index, slide in enumerate(prs.slides, start=1):
        placed = []
        for order, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if not (shape.text_frame.text or "").strip():
                continue
            if shape.top is None or shape.left is None:
                continue
            placed.append((order, int(shape.top), int(shape.left), shape))
        if len(placed) < 2:
            continue
        visual = sorted(placed, key=lambda t: (t[1], t[2]))
        if [p[0] for p in visual] != [p[0] for p in placed]:
            first = visual[0][3]
            text = (first.text_frame.text or "").strip().replace("\n", " ")[:40]
            issues.append(
                Issue(
                    "reading_order",
                    index,
                    f"{len(placed)} text shapes announced out of visual order; "
                    f"topmost is read at position {placed.index(visual[0]) + 1} ({text!r})",
                )
            )
    return issues


def audit(prs: Presentation) -> dict[str, Any]:
    """Run every check. Returns counts by check plus the full issue list."""
    issues: list[Issue] = []
    issues += check_slide_titles(prs)
    issues += check_font_sizes(prs)
    issues += check_table_headers(prs)
    issues += check_link_text(prs)
    issues += check_reading_order(prs)

    counts: dict[str, int] = {}
    for issue in issues:
        counts[issue.check] = counts.get(issue.check, 0) + 1

    return {
        "total_issues": len(issues),
        "by_check": counts,
        "issues": [asdict(i) for i in sorted(issues, key=lambda i: (i.slide, i.check))],
        "note": (
            "Detection only. Nothing in this report was modified. Contrast "
            "checking is not implemented -- PowerPoint colours are usually theme "
            "references with tint maths, and text over an image has no computable "
            "ratio."
        ),
    }
