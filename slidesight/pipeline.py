"""Orchestration: deck in, remediated deck plus report out.

Concurrency is bounded by a semaphore -- firing every image at the gateway at
once gets rate-limited. Results stream through a callback as they land so a UI
can fill a progress bar instead of showing a spinner.
"""

from __future__ import annotations

import asyncio
import time
from pathlib import Path
from typing import Any, Callable

from openai import AsyncOpenAI
from pptx import Presentation

from . import apply, config, wcag
from .describe import describe_image, summarize_deck
from .extract import count_tables, deck_text, extract_images

ProgressFn = Callable[[dict[str, Any]], None]


def make_client() -> AsyncOpenAI:
    """An OpenAI-compatible client pointed at the ASU AIR gateway."""
    return AsyncOpenAI(api_key=config.get_api_key(), base_url=config.get_base_url())


async def _describe_one(
    client: AsyncOpenAI,
    image,
    summary: str,
    semaphore: asyncio.Semaphore,
    on_progress: ProgressFn | None,
) -> dict[str, Any]:
    """Describe one image, converting any failure into a review-queue row."""
    if image.content_type not in config.SUPPORTED_IMAGE_TYPES:
        record = apply.failure_record(
            image, f"unsupported image format {image.content_type}"
        )
    else:
        async with semaphore:
            try:
                result = await describe_image(client, image, summary)
                record = apply.make_record(image, result)
            except Exception as exc:  # noqa: BLE001 - one image must not kill the batch
                record = apply.failure_record(
                    image, f"{type(exc).__name__}: {exc}"[:200]
                )
    if on_progress:
        on_progress(record)
    return record


async def remediate(
    input_path: str | Path,
    output_path: str | Path | None = None,
    *,
    concurrency: int = config.DEFAULT_CONCURRENCY,
    use_summary: bool = True,
    on_progress: ProgressFn | None = None,
) -> dict[str, Any]:
    """Describe every image in a deck, apply what passes the gate, save.

    Returns the report: counts, per-image records, and runtime.
    """
    input_path = Path(input_path)
    started = time.time()

    prs = Presentation(str(input_path))
    images = extract_images(prs)
    tables = count_tables(prs)

    client = make_client()

    summary = ""
    if use_summary and images:
        try:
            summary = await summarize_deck(client, deck_text(prs))
        except Exception:
            summary = ""  # deck summary is a nicety, not a dependency

    semaphore = asyncio.Semaphore(max(1, concurrency))
    records = await asyncio.gather(
        *(
            _describe_one(client, image, summary, semaphore, on_progress)
            for image in images
        ),
        return_exceptions=False,
    )
    records = sorted(records, key=lambda r: (r["slide"], r["image_id"]))

    changed = apply.apply_records(images, records)

    if output_path is not None:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    counts = apply.summarize(records)
    return {
        "source": input_path.name,
        "output": str(output_path) if output_path else None,
        "slides": len(prs.slides),
        "tables_detected": tables,
        "shapes_written": changed,
        "deck_summary": summary,
        "runtime_seconds": round(time.time() - started, 1),
        **counts,
        "images": records,
        "wcag": wcag.audit(prs),
    }


def audit_only(input_path: str | Path) -> dict[str, Any]:
    """WCAG checks with no model calls. Runs in milliseconds.

    This is also the graceful-degradation path: if the vision model is
    unreachable, this alone is still a working accessibility tool.
    """
    input_path = Path(input_path)
    prs = Presentation(str(input_path))
    return {
        "source": input_path.name,
        "slides": len(prs.slides),
        "images_found": len(extract_images(prs)),
        "tables_detected": count_tables(prs),
        "wcag": wcag.audit(prs),
    }
