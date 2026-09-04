"""FastAPI backend wrapping slidesight.remediate.

Endpoints:
    POST /api/upload               -> {"job_id", "status"} (runs pipeline in background)
    GET  /api/jobs/{job_id}        -> live progress / final report / error
    GET  /api/jobs/{job_id}/download -> remediated .pptx as an attachment
    GET  /api/jobs/{job_id}/report -> full report JSON
    GET  /api/health               -> {"status": "ok"}

Design notes:
- remediate() is async (openai AsyncOpenAI), but it must not run on the server's
  own event loop alongside request handling — uvicorn's loop would be starved by
  the model calls' fan-out. Each job therefore gets a worker thread running
  asyncio.run(remediate(...)); the progress callback fires on that thread, so
  the job registry is guarded by a threading.Lock (see jobs.py).
- The progress record from the pipeline has no totals, so a cheap pre-scan
  (open the pptx, count slides/images — no model calls) captures total_slides
  and total_images up front for the progress percentage. If the pre-scan fails
  (e.g. a PDF renamed .pptx), totals stay None and the pipeline's own
  human-readable ValueError becomes the job error.
- Uploaded decks and outputs live in per-job temp dirs under the system temp
  directory, never inside the repo. Jobs older than MAX_JOB_AGE_SECONDS are
  cleaned up (files and registry entry) on the next upload.
"""

from __future__ import annotations

import asyncio
import logging
import os
import shutil
import tempfile
import threading
import uuid
from pathlib import Path

from fastapi import Body, FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, Response
from fastapi.staticfiles import StaticFiles
from pptx import Presentation

from slidesight import remediate
from slidesight.extract import extract_images, set_alt_text

from . import config, slides
from .jobs import registry

logger = logging.getLogger("slidesight.server")
logging.basicConfig(level=logging.INFO)

CHUNK_SIZE = 1024 * 1024  # stream uploads 1 MB at a time

# Reading one thumbnail used to re-open and fully re-parse the deck, loading
# every image blob, on the event loop. A review queue of 15 items meant 15 full
# parses and an unresponsive API. The blobs for a job are read once and kept
# until the job is cleaned up.
_THUMBS: dict[str, dict[tuple[int, str], tuple[bytes, str]]] = {}
_THUMBS_LOCK = threading.Lock()

# Serialises the read-modify-write in approve(), so two approvals cannot
# interleave and corrupt the deck.
_SAVE_LOCKS: dict[str, threading.Lock] = {}

PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

app = FastAPI(title="SlideSight API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=config.ALLOWED_ORIGINS,
    allow_credentials=False,
    allow_methods=["GET", "POST"],
    allow_headers=["*"],
)


def _error(status_code: int, message: str) -> JSONResponse:
    """Consistent error envelope: {"error": "..."} for every failure."""
    return JSONResponse(status_code=status_code, content={"error": message})


def _safe_filename(name: str) -> str:
    """Strip any path components; keep the basename only."""
    return Path(name).name if name else "upload.pptx"


def _prescan_totals(path: Path) -> tuple[int | None, int | None]:
    """Best-effort slide/image counts for the progress bar. No model calls.

    Returns (None, None) if the file cannot be opened here — the pipeline
    re-opens it itself and produces a user-facing error message.
    """
    try:
        prs = Presentation(str(path))
        return len(prs.slides), len(extract_images(prs))
    except Exception:  # noqa: BLE001 - totals are optional
        return None, None


def _run_job(job_id: str, input_path: Path, output_path: Path) -> None:
    """Worker thread target: run the async pipeline to completion."""

    def on_progress(record: dict) -> None:
        # One record per image as it lands:
        # {slide, image_id, alt_text, confidence, decorative, reason, action}
        registry.record_progress(job_id, record)

    try:
        report = asyncio.run(
            remediate(input_path, output_path, on_progress=on_progress)
        )
    except (ValueError, RuntimeError) as exc:
        # ValueError: unreadable/misnamed uploads, message already user-facing.
        # RuntimeError: missing API key, also actionable as-is.
        logger.warning("job %s failed: %s", job_id, type(exc).__name__)
        registry.fail_job(job_id, str(exc))
    except Exception as exc:  # noqa: BLE001 - one job must not kill the server
        logger.exception("job %s failed unexpectedly", job_id)
        registry.fail_job(job_id, f"Remediation failed ({type(exc).__name__}).")
    else:
        registry.complete_job(job_id, report, str(output_path))
        logger.info("job %s complete", job_id)
        # Pictures of the slides, for the accessibility findings. Background, so
        # it never delays the result the user is waiting for.
        slides.start(job_id, input_path, output_path.parent / "slides")


def _save_lock(job_id: str) -> threading.Lock:
    with _THUMBS_LOCK:
        return _SAVE_LOCKS.setdefault(job_id, threading.Lock())


def _thumbs_for(job_id: str, deck_path: Path) -> dict[tuple[int, str], tuple[bytes, str]]:
    """Image bytes for one job, parsed once and reused."""
    with _THUMBS_LOCK:
        cached = _THUMBS.get(job_id)
    if cached is not None:
        return cached
    table = {
        (im.slide, im.image_id): (im.blob, im.content_type)
        for im in extract_images(Presentation(str(deck_path)))
    }
    with _THUMBS_LOCK:
        _THUMBS[job_id] = table
    return table


def _cleanup_old_jobs() -> None:
    """Remove files and registry entries for jobs past MAX_JOB_AGE_SECONDS."""
    for job in registry.jobs_older_than(config.MAX_JOB_AGE_SECONDS):
        job_dir = Path(job["output_path"]).parent if job.get("output_path") else None
        if job_dir and job_dir.is_relative_to(config.UPLOAD_ROOT):
            shutil.rmtree(job_dir, ignore_errors=True)
        registry.drop_job(job["job_id"])
        with _THUMBS_LOCK:
            _THUMBS.pop(job["job_id"], None)
            _SAVE_LOCKS.pop(job["job_id"], None)
        slides.forget(job["job_id"])
        logger.info("cleaned up expired job %s", job["job_id"])


@app.get("/api/health")
async def health() -> dict:
    return {"status": "ok"}


@app.post("/api/upload")
async def upload(file: UploadFile = File(...)) -> JSONResponse:
    name = _safe_filename(file.filename or "")
    lower = name.lower()

    # Reject out-of-scope formats on the extension, before any pipeline work,
    # with the same messages the CLI uses.
    if lower.endswith(".pdf"):
        return _error(
            400,
            f"{name} is a PDF. SlideSight reads PowerPoint files only — for PDFs "
            "see ASU CIC's PDF Accessibility tool.",
        )
    if lower.endswith(".ppt"):
        return _error(
            400,
            f"{name} is a legacy .ppt. Convert it first: "
            "soffice --headless --convert-to pptx",
        )
    if not lower.endswith(".pptx"):
        return _error(400, f"{name} is not a .pptx file. SlideSight reads PowerPoint files only.")

    _cleanup_old_jobs()

    job_id = uuid.uuid4().hex
    job_dir = config.UPLOAD_ROOT / f"job_{job_id}"
    job_dir.mkdir(parents=True, exist_ok=True)

    # Stream to disk without holding the whole deck in memory.
    input_path = job_dir / name
    try:
        with input_path.open("wb") as out:
            while chunk := await file.read(CHUNK_SIZE):
                out.write(chunk)
    except OSError as exc:
        shutil.rmtree(job_dir, ignore_errors=True)
        logger.error("upload write failed: %s", type(exc).__name__)
        return _error(500, "Could not save the uploaded file. Try again.")

    total_slides, total_images = _prescan_totals(input_path)
    output_path = job_dir / f"remediated_{name}"
    registry.create_job(
        name, total_slides=total_slides, total_images=total_images, job_id=job_id
    )

    worker = threading.Thread(
        target=_run_job, args=(job_id, input_path, output_path),
        name=f"remediate-{job_id[:8]}", daemon=True,
    )
    worker.start()

    return JSONResponse(status_code=202, content={"job_id": job_id, "status": "processing"})


def _processing_state(job: dict) -> dict:
    state = {
        "job_id": job["job_id"],
        "status": job["status"],
        "progress_pct": job["progress_pct"],
        "current_slide": job["current_slide"],
        "total_slides": job["total_slides"],
        "total_images": job["total_images"],
        "current_image_id": job["current_image_id"],
        "current_alt_text": job["current_alt_text"],
        "current_confidence": job["current_confidence"],
        "current_action": job["current_action"],
        "records_so_far": job["records_so_far"],
    }
    if job.get("error"):
        state["error"] = job["error"]
    return state


@app.get("/api/jobs/{job_id}")
async def job_state(job_id: str):
    job = registry.get_job(job_id)
    if job is None:
        return _error(404, f"Unknown job: {job_id}")

    if job["status"] == "complete":
        # Never leak server filesystem paths to the client.
        job.pop("output_path", None)
        # Sanitize any report output to replace temp paths with just the filename
        if job.get("report"):
            report = dict(job["report"])
            report["output"] = job.get("output_filename") or report.get("output", "")
            job["report"] = report
        return job
    return _processing_state(job)


@app.get("/api/jobs/{job_id}/download")
async def download(job_id: str):
    job = registry.get_job(job_id)
    if job is None:
        return _error(404, f"Unknown job: {job_id}")
    if job["status"] != "complete":
        return _error(404, f"Job {job_id} is {job['status']}; no file to download yet.")

    output_path = Path(job["output_path"])
    if not output_path.is_file():
        return _error(404, f"Remediated file for job {job_id} is missing.")
    return FileResponse(
        output_path,
        media_type=PPTX_MEDIA_TYPE,
        filename=job["output_filename"],
    )


@app.get("/api/jobs/{job_id}/report")
async def report(job_id: str):
    job = registry.get_job(job_id)
    if job is None:
        return _error(404, f"Unknown job: {job_id}")
    if job["status"] != "complete" or job["report"] is None:
        return _error(404, f"Job {job_id} is {job['status']}; no report yet.")
    # Never leak server filesystem paths to the client.
    report = dict(job["report"])
    report["output"] = job.get("output_filename") or report.get("output", "")
    return report


def _find_image(prs, slide_no: int, image_id: str):
    """Locate one picture by (slide, image_id) in an open Presentation."""
    slides = list(prs.slides)
    if not 1 <= slide_no <= len(slides):
        return None
    for image in extract_images(prs):
        if image.slide == slide_no and image.image_id == image_id:
            return image
    return None


@app.get("/api/jobs/{job_id}/thumb/{slide_no}/{image_id}")
def thumbnail(job_id: str, slide_no: int, image_id: str, w: int = 480):
    """The actual picture bytes, so the review queue can show what it is asking about.

    A reviewer cannot judge a description without seeing the image. Served from
    the job's own deck, downscaled, and never cached beyond the job's lifetime.
    """
    job = registry.get_job(job_id)
    if job is None:
        return _error(404, f"Unknown job: {job_id}")

    source = job.get("output_path") or ""
    if not source or not Path(source).is_file():
        return _error(404, "No deck available for this job yet.")

    try:
        found = _thumbs_for(job_id, Path(source)).get((slide_no, image_id))
    except Exception:  # noqa: BLE001
        return _error(404, "Could not read the deck.")
    if found is None:
        return _error(404, f"No image {image_id} on slide {slide_no}.")

    from slidesight.describe import prepare_image

    blob, content_type = found
    # The review list needs a small tile; the preview overlay needs something a
    # person can actually read a chart from. Same picture, two sizes.
    max_edge = max(120, min(int(w), 1600))
    prepared = prepare_image(blob, content_type, max_edge=max_edge)
    if prepared is None:
        return _error(415, "This image format cannot be displayed.")
    blob, content_type = prepared
    return Response(content=blob, media_type=content_type,
                    headers={"Cache-Control": "private, max-age=3600"})


@app.get("/api/jobs/{job_id}/slide/{slide_no}")
def slide_image(job_id: str, slide_no: int):
    """A picture of one slide, for the accessibility findings.

    404 while the background render is still running, or if LibreOffice is not
    available. The page hides the image rather than showing a broken one.
    """
    if registry.get_job(job_id) is None:
        return _error(404, f"Unknown job: {job_id}")
    path = slides.slide_png(job_id, slide_no)
    if path is None:
        return _error(404, "No picture for that slide yet.")
    return FileResponse(path, media_type="image/png",
                        headers={"Cache-Control": "private, max-age=3600"})


@app.post("/api/jobs/{job_id}/approve")
async def approve(job_id: str, payload: dict = Body(...)):
    """Write a human-approved description into the remediated deck.

    This is the point of the review queue: the pipeline deliberately did NOT
    write these, so approving is what puts them in the file. Editing the draft
    first is the normal path -- the text sent here is what gets written, not the
    model's original.
    """
    job = registry.get_job(job_id)
    if job is None:
        return _error(404, f"Unknown job: {job_id}")
    if job["status"] != "complete":
        return _error(409, f"Job {job_id} is {job['status']}; nothing to write yet.")

    try:
        slide_no = int(payload.get("slide"))
        image_id = str(payload.get("image_id"))
    except (TypeError, ValueError):
        return _error(400, "slide and image_id are required.")
    alt_text = str(payload.get("alt_text") or "").strip()
    if not alt_text:
        return _error(400, "alt_text cannot be empty. Use skip to leave it unwritten.")

    output_path = Path(job["output_path"])
    if not output_path.is_file():
        return _error(404, "The remediated file is missing.")

    try:
        with _save_lock(job_id):
            prs = Presentation(str(output_path))
            image = _find_image(prs, slide_no, image_id)
            if image is None:
                return _error(404, f"No image {image_id} on slide {slide_no}.")
            set_alt_text(image.shape, alt_text)
            # Save beside the target and swap it in. Writing a zip in place
            # truncates it first, so an interrupted save would leave a
            # half-written file that still passes the is_file() check on
            # download and will not open in PowerPoint.
            tmp = output_path.with_name(output_path.name + ".tmp")
            prs.save(str(tmp))
            os.replace(tmp, output_path)
    except Exception as exc:  # noqa: BLE001
        logger.exception("approve failed for job %s", job_id)
        return _error(500, f"Could not write the description ({type(exc).__name__}).")

    # Keep the stored report in step with the file the user will download.
    # Rebuild the report so the JSON a user downloads agrees with the .pptx.
    # get_job returns report by reference, so copy before mutating.
    from slidesight import apply as _apply

    source_report = job.get("report") or {}
    report = dict(source_report)
    images = [dict(r) for r in report.get("images", [])]
    for record in images:
        if record["slide"] == slide_no and record["image_id"] == image_id:
            record["alt_text"] = alt_text
            record["action"] = "human_approved"
            break
    report["images"] = images
    report.update(_apply.summarize(images))
    report["human_approved"] = sum(1 for r in images if r["action"] == "human_approved")
    registry.update_job(job_id, report=report)

    logger.info("job %s: approved %s on slide %s", job_id, image_id, slide_no)
    return {"status": "written", "slide": slide_no, "image_id": image_id}


@app.middleware("http")
async def revalidate_web_assets(request, call_next):
    """Make the browser check for a new app.js instead of reusing a stale one.

    StaticFiles sends an ETag but no Cache-Control, so a browser is free to
    serve the previous copy from its heuristic cache. During a demo that means
    editing the UI appears to do nothing. no-cache still allows a 304, so this
    costs a round trip, not a download.
    """
    response = await call_next(request)
    if not request.url.path.startswith("/api/"):
        response.headers["Cache-Control"] = "no-cache"
    return response


# The web UI is served from this same app, so one command runs everything.
WEB_ROOT = Path(__file__).resolve().parent.parent / "web"
if WEB_ROOT.is_dir():
    app.mount("/", StaticFiles(directory=str(WEB_ROOT), html=True), name="web")


if __name__ == "__main__":
    import uvicorn

    uvicorn.run("server.main:app", host="0.0.0.0", port=8000)
