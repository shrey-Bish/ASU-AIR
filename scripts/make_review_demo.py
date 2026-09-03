"""Build a demo deck that deliberately trips the confidence gate.

    .venv/bin/python scripts/make_review_demo.py "decks/some.pptx" -o data/demo/

Real lecture decks are mostly clean, so the review queue fires on roughly 1
image in 500 -- which is honest, but leaves nothing to show. This takes a real
deck and degrades exactly one image in a measured way (Gaussian blur), so the
gate has something to catch.

**This is a constructed demo, not a found failure.** The degradation is ours.
Say so when showing it. What it demonstrates is real: the same pipeline, on the
same deck, routing a damaged image to a human instead of inventing a
description for it.

The original image is kept alongside so you can show before and after.
"""

from __future__ import annotations

import argparse
import io
import shutil
import zipfile
from pathlib import Path

from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import sys

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from slidesight.extract import walk  # noqa: E402


def find_target(path: Path) -> tuple[int, str, str]:
    """Pick the image to degrade: the largest picture on the deck.

    Largest is a good proxy for "the one carrying the content", which makes the
    before/after comparison mean something.
    """
    prs = Presentation(str(path))
    best = None
    for index, slide in enumerate(prs.slides, start=1):
        for shape in walk(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
                part = shape.image  # noqa: F841 - accessed for the exception check
            except Exception:
                continue
            rid = shape._element.blip_rId
            image_part = shape.part.related_part(rid)
            name = str(image_part.partname).lstrip("/")
            if best is None or len(blob) > best[3]:
                best = (index, name, shape.name, len(blob))
    if best is None:
        raise SystemExit(f"no usable images in {path.name}")
    return best[0], best[1], best[2]


def blur_entry(src: Path, dst: Path, member: str, radius: float) -> tuple[bytes, bytes]:
    """Rewrite the pptx zip with one media file replaced by a blurred copy.

    Editing the zip directly keeps every other part byte-identical, so the only
    difference between the two decks is the one image.
    """
    with zipfile.ZipFile(src) as zin:
        original = zin.read(member)
        img = Image.open(io.BytesIO(original))
        fmt = img.format or "PNG"
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
            fmt = "JPEG"
        blurred_img = img.filter(ImageFilter.GaussianBlur(radius))
        buf = io.BytesIO()
        blurred_img.save(buf, format=fmt)
        blurred = buf.getvalue()

        dst.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = blurred if item.filename == member else zin.read(item.filename)
                zout.writestr(item, data)
    return original, blurred


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck", help="a real .pptx to base the demo on")
    ap.add_argument("-o", "--out", default="data/demo", help="output folder")
    ap.add_argument("-r", "--radius", type=float, default=12.0, help="blur radius")
    args = ap.parse_args()

    src = Path(args.deck)
    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    slide_no, member, shape_name = find_target(src)
    demo = out_dir / f"{src.stem}__review_demo.pptx"
    original, blurred = blur_entry(src, demo, member, args.radius)

    clean = out_dir / f"{src.stem}__original.pptx"
    shutil.copy2(src, clean)

    (out_dir / "before.png").write_bytes(original)
    (out_dir / "after_blurred.png").write_bytes(blurred)

    print(f"target      slide {slide_no}, shape {shape_name!r} ({member})")
    print(f"blur radius {args.radius}")
    print(f"clean deck  {clean}")
    print(f"demo deck   {demo}")
    print(f"images      {out_dir/'before.png'} / {out_dir/'after_blurred.png'}")
    print(
        "\nRun the pipeline on both. The clean deck should describe this image;\n"
        "the demo deck should send it to the review queue with a stated reason."
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
