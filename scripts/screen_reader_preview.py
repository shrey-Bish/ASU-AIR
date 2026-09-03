"""Hear what a screen reader announces, before and after remediation.

    .venv/bin/python scripts/screen_reader_preview.py ORIGINAL.pptx REMEDIATED.pptx
    .venv/bin/python scripts/screen_reader_preview.py orig.pptx fixed.pptx --slide 13
    .venv/bin/python scripts/screen_reader_preview.py orig.pptx fixed.pptx --slide 13 --save demo

Uses the macOS `say` voice, which is the same speech engine VoiceOver uses.
This is not a substitute for running VoiceOver -- it does not prove PowerPoint
exposes the field -- but it is exactly what a blind student hears for each
image, and `--save` writes audio files you can drop straight into the video.

A picture with no alt text is announced by its shape name ("Picture 4"). That
is the "before".
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from slidesight.extract import walk  # noqa: E402

VOICE = "Samantha"


def announcements(path: Path, slide_no: int | None) -> list[tuple[int, str, str]]:
    """What a screen reader says for each picture: (slide, shape name, spoken)."""
    prs = Presentation(str(path))
    out = []
    for index, slide in enumerate(prs.slides, start=1):
        if slide_no and index != slide_no:
            continue
        for shape in walk(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            attrib = shape._element._nvXxPr.cNvPr.attrib
            alt = (attrib.get("descr") or "").strip()
            if alt:
                spoken = alt
            elif "descr" in attrib:
                spoken = "(skipped — marked decorative)"
            else:
                # No alt text at all: the reader falls back to the shape name.
                spoken = attrib.get("name") or shape.name or "image"
            out.append((index, shape.name, spoken))
    return out


def speak(text: str, save: Path | None) -> None:
    if save:
        subprocess.run(["say", "-v", VOICE, "-o", str(save), text], check=False)
        print(f"      saved {save}")
    else:
        subprocess.run(["say", "-v", VOICE, text], check=False)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("original")
    ap.add_argument("remediated")
    ap.add_argument("--slide", type=int, help="only this slide number")
    ap.add_argument("--save", help="write audio files with this prefix instead of playing")
    ap.add_argument("--silent", action="store_true", help="print only, do not speak")
    args = ap.parse_args()

    before = announcements(Path(args.original), args.slide)
    after = announcements(Path(args.remediated), args.slide)

    if not before:
        print("no pictures found (check the slide number)")
        return 1

    for i, ((sl, name, b), (_, _, a)) in enumerate(zip(before, after)):
        print(f"\n--- slide {sl}, {name}")
        print(f"  BEFORE: {b}")
        print(f"  AFTER : {a}")
        if args.silent:
            continue
        if args.save:
            speak(b, Path(f"{args.save}_{sl}_{i}_before.aiff"))
            speak(a, Path(f"{args.save}_{sl}_{i}_after.aiff"))
        else:
            print("  (speaking before…)")
            speak(b, None)
            print("  (speaking after…)")
            speak(a, None)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
